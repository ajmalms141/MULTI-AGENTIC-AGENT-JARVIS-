from fastapi import FastAPI
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI(title="Jarvis Presentation Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

import os
import re
import json
import httpx
from fastapi import HTTPException
from fastapi.responses import FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Output folder for generated presentations
OUTPUT_DIR = "/output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

@app.get("/health")
def health():
    return {"status": "ok", "service": "presentation"}

def get_gemini_keys() -> list:
    """Collect all configured Gemini API keys (GEMINI_API_KEY, GEMINI_API_KEY_2, GEMINI_API_KEY_3, ...)."""
    keys = []
    # Always check base key
    base = os.environ.get("GEMINI_API_KEY", "").strip()
    if base:
        keys.append(base)
    # Check numbered extras
    for i in range(2, 10):
        extra = os.environ.get(f"GEMINI_API_KEY_{i}", "").strip()
        if extra:
            keys.append(extra)
        else:
            break
    return keys

async def call_gemini_with_rotation(payload: dict) -> dict:
    """Try each Gemini API key in order, rotating on rate limit (429) or errors."""
    keys = get_gemini_keys()
    if not keys:
        raise HTTPException(status_code=500, detail="No Gemini API key configured on server")
    
    last_error = None
    for idx, key in enumerate(keys):
        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={key}"
        try:
            async with httpx.AsyncClient() as client:
                response = await client.post(url, json=payload, timeout=40.0)
            
            if response.status_code == 200:
                print(f"Gemini call succeeded with key #{idx + 1}")
                return response.json()
            elif response.status_code == 429:
                print(f"Key #{idx + 1} rate limited (429), trying next key...")
                last_error = f"Key #{idx + 1}: Rate limited (429)"
                continue
            else:
                print(f"Key #{idx + 1} returned error {response.status_code}: {response.text[:200]}")
                last_error = f"Key #{idx + 1}: HTTP {response.status_code}"
                continue
        except Exception as e:
            print(f"Key #{idx + 1} request failed: {e}")
            last_error = str(e)
            continue
    
    raise HTTPException(status_code=503, detail=f"All Gemini API keys failed. Last error: {last_error}")

@app.get("/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename=filename)

async def search_unsplash_image(query: str) -> str:
    access_key = os.environ.get("UNSPLASH_ACCESS_KEY")
    if not access_key:
        return ""
    try:
        url = f"https://api.unsplash.com/search/photos?query={httpx.QueryParams(query)}&client_id={access_key}&per_page=1"
        async with httpx.AsyncClient() as client:
            response = await client.get(url, timeout=5.0)
            if response.status_code == 200:
                data = response.json()
                if data.get("results"):
                    return data["results"][0]["urls"]["regular"]
    except Exception as e:
        print(f"Unsplash search error for query '{query}': {e}")
    return ""

async def download_image(url: str, filepath: str) -> bool:
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(url, timeout=10.0)
            if response.status_code == 200:
                with open(filepath, "wb") as f:
                    f.write(response.content)
                return True
    except Exception as e:
        print(f"Image download error for URL '{url}': {e}")
    return False

@app.post("/generate")
async def generate(body: dict):
    topic = body.get("topic", "").strip()
    slides_count = int(body.get("slides_count", 5))
    
    if not topic:
        raise HTTPException(status_code=400, detail="Topic is required")
        
    # Validate at least one key is available
    if not get_gemini_keys():
        raise HTTPException(status_code=500, detail="No Gemini API key configured on server")
        
    # Query Gemini to generate slide structure (with automatic key rotation)
    system_prompt = (
        "You are an expert presentation designer. Create a highly informative presentation outline in JSON format. "
        "Each slide must contain a slide 'title', an array of 3-4 bullet 'points', and a visual 'image_query' keyword "
        "to search on Unsplash (2-3 search terms max). Format the output as JSON conforming strictly to this structure: "
        '{"slides": [{"title": "...", "points": ["...", "..."], "image_query": "..."}]}. '
        "Do not output markdown code blocks. Just return the JSON object."
    )
    user_prompt = f"Create a presentation about: '{topic}' containing exactly {slides_count} slides."
    
    payload = {
        "contents": [{
            "parts": [{"text": f"{system_prompt}\n\n{user_prompt}"}]
        }],
        "generationConfig": {
            "responseMimeType": "application/json"
        }
    }
    
    try:
        # call_gemini_with_rotation tries each key automatically on 429/error
        raw_response = await call_gemini_with_rotation(payload)
        content_text = raw_response["candidates"][0]["content"]["parts"][0]["text"]
        
        # Clean up the output string to ensure it's valid JSON
        clean_json_str = content_text.strip()
        if clean_json_str.startswith("```json"):
            clean_json_str = clean_json_str.replace("```json", "", 1)
        if clean_json_str.endswith("```"):
            clean_json_str = clean_json_str[:-3]
        clean_json_str = clean_json_str.strip()
        
        slides_data = json.loads(clean_json_str)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate content with Gemini: {str(e)}")
        
    # Build presentation using python-pptx
    prs = Presentation()
    # Use widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Premium Color Scheme (Dark Slate Theme)
    BG_COLOR = RGBColor(11, 15, 25)       # Deep slate navy
    PRIMARY_COLOR = RGBColor(255, 255, 255) # White
    ACCENT_COLOR = RGBColor(59, 130, 246)  # Neon blue
    TEXT_COLOR = RGBColor(209, 213, 219)   # Light gray
    
    slides = slides_data.get("slides", [])
    
    for index, slide_info in enumerate(slides):
        # Add a blank slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Apply dark background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
        # Unsplash Image Search & Download
        image_query = slide_info.get("image_query", topic)
        image_url = await search_unsplash_image(image_query)
        image_path = None
        
        if image_url:
            temp_img_name = f"temp_{index}_{uuid_str()}.jpg"
            temp_img_path = os.path.join(OUTPUT_DIR, temp_img_name)
            if await download_image(image_url, temp_img_path):
                image_path = temp_img_path

        # Title formatting
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_info.get("title", f"Slide {index + 1}")
        p.font.name = "Arial"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR
        
        # Add decorative accent line
        line = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(0.8), Inches(1.5), Inches(2.0), Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.color.rgb = ACCENT_COLOR
        
        # Layout splits (Text left, Image right)
        if image_path:
            # Bullet points left box
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5))
            tf_bullets = text_box.text_frame
            tf_bullets.word_wrap = True
            
            for pt_idx, point_text in enumerate(slide_info.get("points", [])):
                p_bullet = tf_bullets.add_paragraph() if pt_idx > 0 else tf_bullets.paragraphs[0]
                p_bullet.text = f"•  {point_text}"
                p_bullet.font.name = "Arial"
                p_bullet.font.size = Pt(18)
                p_bullet.font.color.rgb = TEXT_COLOR
                p_bullet.space_after = Pt(14)
                
            # Image right box
            try:
                # Add picture on right side: left=7.3", top=2.0", width=5.2", height=4.5"
                slide.shapes.add_picture(image_path, Inches(7.3), Inches(2.0), Inches(5.2), Inches(4.5))
            except Exception as pic_err:
                print(f"Failed to place picture in slide: {pic_err}")
                
            # Clean up temp image
            try:
                os.remove(image_path)
            except:
                pass
        else:
            # Full width text box if no image
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.5))
            tf_bullets = text_box.text_frame
            tf_bullets.word_wrap = True
            
            for pt_idx, point_text in enumerate(slide_info.get("points", [])):
                p_bullet = tf_bullets.add_paragraph() if pt_idx > 0 else tf_bullets.paragraphs[0]
                p_bullet.text = f"•  {point_text}"
                p_bullet.font.name = "Arial"
                p_bullet.font.size = Pt(22)
                p_bullet.font.color.rgb = TEXT_COLOR
                p_bullet.space_after = Pt(18)
                
        # Footer slide number
        slide_num_box = slide.shapes.add_textbox(Inches(11.7), Inches(6.8), Inches(0.8), Inches(0.4))
        p_num = slide_num_box.text_frame.paragraphs[0]
        p_num.text = str(index + 1)
        p_num.font.name = "Arial"
        p_num.font.size = Pt(12)
        p_num.font.color.rgb = TEXT_COLOR
        
    # Save the PPTX
    topic_slug = re.sub(r'[^a-zA-Z0-9]', '_', topic.lower())
    filename = f"{topic_slug}_{uuid_str()[:8]}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    
    try:
        prs.save(filepath)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save PPTX file: {str(e)}")
        
    return {
        "status": "success",
        "topic": topic,
        "filename": filename,
        "download_url": f"/download/{filename}"
    }

def uuid_str() -> str:
    import uuid
    return str(uuid.uuid4())

